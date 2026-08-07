"""Outils de l'agent codeur.

Un outil, c'est deux choses :
  1. Un SCHÉMA JSON  -> décrit l'outil au modèle (nom, paramètres, utilité).
  2. Une FONCTION     -> fait réellement le travail quand le modèle l'appelle.

Le modèle ne "voit" que les schémas. Il décide d'appeler un outil en
renvoyant un JSON. Notre code exécute la fonction correspondante et
renvoie le résultat textuel au modèle. C'est ça, un "tool" d'agent.
"""

import glob as globlib
import math
import os
import re
import shutil
import subprocess
import time
import zipfile
import zlib

import config


# --------------------------------------------------------------------------
# 1) Schémas JSON des outils (format "function calling" de OpenAI/Ollama)
# --------------------------------------------------------------------------

TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "list_dir",
            "description": "Liste le contenu d'un répertoire (fichiers + dossiers). "
                           "À utiliser pour explorer un projet avant de travailler.",
            "parameters": {
                "type": "object",
                "properties": {"path": {"type": "string", "description": "Chemin du dossier (défaut: .)"}},
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": "Lit le contenu d'un fichier, lignes numérotées. "
                           "TOUJOURS lire un fichier avant de le modifier.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Chemin du fichier"},
                    "offset": {"type": "integer", "description": "Ligne de départ (1 = début)"},
                    "limit": {"type": "integer", "description": "Nombre de lignes à lire"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_document",
            "description": "Extrait le TEXTE lisible d'un document : Word (.docx), "
                           "PowerPoint (.pptx), Excel (.xlsx), PDF (.pdf) ou fichier texte. "
                           "À utiliser pour tout fichier qui n'est pas du code simple.",
            "parameters": {
                "type": "object",
                "properties": {"path": {"type": "string", "description": "Chemin du fichier document"}},
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "search_in_files",
            "description": "Recherche un mot ou une expression dans tous les fichiers texte "
                           "d'un dossier (récursif). Retourne les fichiers et les NUMÉROS DE "
                           "LIGNE contenant l'expression. Idéal pour 'où est utilisé X ?'.",
            "parameters": {
                "type": "object",
                "properties": {
                    "folder": {"type": "string", "description": "Dossier à parcourir (défaut: .)"},
                    "term": {"type": "string", "description": "Mot ou expression à chercher"},
                },
                "required": ["term"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "count_occurrences",
            "description": "Compte EXACTEMENT combien de fois un mot ou une expression apparaît "
                           "dans un fichier (décompte précis par le code, pas à la main). "
                           "Compte les mots entiers, sans tenir compte des majuscules.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Chemin du fichier à analyser"},
                    "term": {"type": "string", "description": "Mot ou expression à compter"},
                },
                "required": ["path", "term"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_file",
            "description": "Crée un fichier ou ÉCRASE son contenu complet avec le texte fourni. "
                           "À utiliser pour créer de nouveaux fichiers.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Chemin du fichier"},
                    "content": {"type": "string", "description": "Contenu complet à écrire"},
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "edit_file",
            "description": "Remplace une portion EXACTE d'un fichier par un nouveau texte. "
                           "Pour modifier sans réécrire tout le fichier. "
                           "Échoue si l'ancien texte n'est pas trouvé à l'identique.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Chemin du fichier"},
                    "old_string": {"type": "string", "description": "Texte EXACT à remplacer"},
                    "new_string": {"type": "string", "description": "Nouveau texte"},
                },
                "required": ["path", "old_string", "new_string"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "glob",
            "description": "Trouve des fichiers par motif de nom, ex: '**/*.py' ou 'src/**'. "
                           "Comme la commande glob de VS Code.",
            "parameters": {
                "type": "object",
                "properties": {"pattern": {"type": "string", "description": "Motif à rechercher"}},
                "required": ["pattern"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_image",
            "description": "Décrit une IMAGE (jpg, png, bmp, webp...) : ce qu'elle "
                           "montre et le texte visible. À utiliser pour TOUTE "
                           "image jointe ou tout chemin d'image. Fonctionne hors "
                           "ligne (modèle de vision local ou OCR).",
            "parameters": {
                "type": "object",
                "properties": {"path": {"type": "string", "description": "Chemin du fichier image"}},
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "generer_image",
            "description": "Génère une IMAGE (PNG) à partir d'une description, 100 % "
                           "hors ligne. Ordre des moteurs : (1) Stable Diffusion déjà "
                           "lancé sur le poste (AGENT_SD_URL) ; (2) stable-diffusion.cpp "
                           "avec son modèle (AGENT_SDCPP + AGENT_SD_MODEL) ; (3) sinon une "
                           "ILLUSTRATION locale (graphiques en barres/secteurs, "
                           "organigrammes, mind-maps, tableaux de comparaison, formules). "
                           "Pour un graphique, fournissez les données en 'label=valeur' "
                           "séparés par ';' (ex: '2020=80 ; 2021=90').",
            "parameters": {
                "type": "object",
                "properties": {
                    "prompt": {"type": "string",
                               "description": "Description de l'image à générer, ou "
                                              "données du graphique en 'label=valeur'"},
                    "sortie": {"type": "string",
                               "description": "Chemin du fichier PNG à créer "
                                              "(défaut: dossier images/ de l'espace de travail)"},
                },
                "required": ["prompt"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "creer_powerpoint",
            "description": "Crée une présentation PowerPoint (.pptx) de cours pour "
                           "enseignant à partir d'un PLAN en JSON : diapositives avec "
                           "titres, puces, notes, IMAGES (chemins de fichiers PNG/JPG "
                           "existants) et LIENS VIDÉO cliquables. Ne JAMAIS inventer "
                           "d'URL de vidéo : n'utiliser que des liens réels fournis.",
            "parameters": {
                "type": "object",
                "properties": {
                    "plan": {"type": "string",
                             "description": "Plan JSON de la présentation, ex: "
                                            "{\"titre\": \"Cours de maths\", \"auteur\": \"M. Dupont\", "
                                            "\"slides\": [{\"titre\": \"Objectifs\", \"texte\": "
                                            "[\"puce 1\", \"puce 2\"], \"notes\": \"À dire en classe\", "
                                            "\"image\": \"chemin.png\", \"video\": {\"url\": "
                                            "\"https://...\", \"texte\": \"Voir la vidéo\"}}]}"},
                    "sortie": {"type": "string",
                               "description": "Chemin du fichier .pptx à créer (défaut: "
                                              "dossier presentations/ de l'espace de travail)"},
                },
                "required": ["plan"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "calcul_symbolique",
            "description": "Calcule et VÉRIFIE une expression mathématique avec SymPy "
                           "(intégrale, dérivée, équation, limite, simplification). "
                           "À utiliser pour TOUT calcul mathématique. L'expression doit être "
                           "copiée EN TEXTE BRUT telle qu'écrite par l'étudiant, ex: "
                           "'ln(x+1)', 'x^2 - 5*x + 6 = 0'. Ne jamais écrire le code SymPy "
                           "à la main : c'est l'outil qui interprète et vérifie.",
            "parameters": {
                "type": "object",
                "properties": {
                    "expression": {"type": "string",
                                   "description": "La fonction mathématique EN TEXTE BRUT, copiée telle quelle (ex: 'ln(x+1)' ou '∫ln(x+1)dx' — l'outil gère le symbole ∫ et le dx)"},
                    "operation": {"type": "string",
                                  "description": "Type de calcul: 'integrale' (défaut), 'derivee', 'equation', 'limite', 'simplifier', 'factoriser'"},
                    "variable": {"type": "string", "description": "Variable d'intégration (défaut: x)"},
                    "point": {"type": "string", "description": "Point pour une limite (défaut: oo)"},
                },
                "required": ["expression"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "bash",
            "description": "Exécute une commande dans le terminal (python, git, pip, pytest...). "
                           "Retourne la sortie, les erreurs et le code de sortie. "
                           "Pour un CALCUL MATHÉMATIQUE (intégrale, dérivée, équation...), "
                           "n'écris PAS le code SymPy à la main : utilise l'outil dédié "
                           "`calcul_symbolique` avec l'expression en texte brut.",
            "parameters": {
                "type": "object",
                "properties": {
                    "command": {"type": "string", "description": "Commande shell à exécuter"},
                    "timeout": {"type": "integer", "description": "Délai max en secondes"},
                },
                "required": ["command"],
            },
        },
    },
]


# --------------------------------------------------------------------------
# 2) Implémentations Python des outils
# --------------------------------------------------------------------------

def _repair_path(path: str) -> str:
    """Corrige un chemin à racine dupliquée (erreur de copier-coller).

    Ex: 'C:\\LaSalle\\C:\\LaSalle\\f.pdf' -> 'C:\\LaSalle\\f.pdf'
    On garde la racine la plus à droite (le dernier 'X:\\').
    """
    if not path:
        return path
    norm = path.replace("/", "\\")
    roots = re.findall(r"[A-Za-z]:\\", norm)
    if len(roots) >= 2:
        return norm[norm.rfind(roots[-1]):]
    return path


def _existing_path(path: str) -> str:
    """Renvoie un chemin existant (réparé si besoin), sinon le chemin d'origine."""
    if os.path.exists(path):
        return path
    repaired = _repair_path(path)
    return repaired if os.path.exists(repaired) else path


def list_dir(path: str = ".") -> str:
    path = _existing_path(path)
    if not os.path.isdir(path):
        return f"ERREUR: '{path}' n'est pas un dossier."
    rows = []
    for name in sorted(os.listdir(path)):
        full = os.path.join(path, name)
        kind = "📁" if os.path.isdir(full) else "📄"
        size = "" if os.path.isdir(full) else f"  {os.path.getsize(full):,} o"
        rows.append(f"{kind}  {name}{size}")
    return "\n".join(rows) or "(dossier vide)"


def read_file(path: str, offset: int = 1, limit: int | None = None) -> str:
    path = _existing_path(path)
    if not os.path.isfile(path):
        parent = os.path.dirname(os.path.abspath(path))
        try:
            listing = "\n".join(
                f"  - {name}" for name in sorted(os.listdir(parent))[:40])
        except OSError:
            listing = ""
        hint = (f"\nListe du dossier {parent}:\n{listing}" if listing else "")
        return (f"ERREUR: fichier introuvable: {path}"
                f"\nVérifiez le chemin ou cherchez le fichier.{hint}")
    limit = limit or config.MAX_READ_LINES
    try:
        with open(path, "rb") as fh:
            data = fh.read()
    except OSError as err:
        return f"ERREUR: {err}"
    if len(data) > 5 * 1024 * 1024:                  # on protège le contexte du modèle
        return f"ERREUR: fichier trop volumineux ({len(data)} octets)."
    if b"\0" in data[:8192]:                         # un .docx (binaire ZIP) contient des NUL
        return (f"ERREUR: '{path}' est un fichier BINAIRE (non texte). "
                f"Pour un document, utilisez l'outil read_document.")
    if os.path.splitext(path)[1].lower() in IMAGE_EXTS:
        return (f"ERREUR: '{path}' est une IMAGE, elle ne se lit pas avec "
                f"read_file. Utilisez l'outil read_image pour la décrire.")
    lines = data.decode("utf-8", errors="replace").split("\n")
    chunk = lines[offset - 1: offset - 1 + limit]
    return "".join(f"{offset + i}: {line}\n" for i, line in enumerate(chunk))


def write_file(path: str, content: str) -> str:
    os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
    with open(path, "w", encoding="utf-8") as fh:
        fh.write(content)
    return f"OK: {path} écrit ({len(content)} caractères)."


def edit_file(path: str, old_string: str, new_string: str) -> str:
    if not os.path.isfile(path):
        return f"ERREUR: fichier introuvable: {path}"
    with open(path, encoding="utf-8") as fh:
        content = fh.read()
    count = content.count(old_string)
    if count == 0:
        return f"ERREUR: texte à remplacer INTROUVABLE dans {path}. Relisez le fichier."
    if count > 1:
        return f"ERREUR: texte trouvé {count} fois. Donnez plus de contexte."
    with open(path, "w", encoding="utf-8") as fh:
        fh.write(content.replace(old_string, new_string, 1))
    return f"OK: {path} modifié (1 remplacement)."


def glob(pattern: str) -> str:
    matches = globlib.glob(pattern, recursive=True)
    return "\n".join(matches) if matches else f"Aucun fichier pour '{pattern}'."


def bash(command: str, timeout: int | None = None) -> str:
    timeout = timeout or config.BASH_TIMEOUT
    started = time.monotonic()
    try:
        proc = subprocess.run(
            command, shell=True, capture_output=True, text=True,
            timeout=timeout, encoding="utf-8", errors="replace",
        )
    except subprocess.TimeoutExpired:
        return f"ERREUR: commande interrompue après {timeout} s."
    elapsed = time.monotonic() - started
    output = proc.stdout + proc.stderr
    if len(output) > 4000:                       # tronque pour ne pas saturer le contexte
        output = output[:4000] + "\n...[sortie tronquée]"
    return (f"$ {command}  ({elapsed:.1f}s, code de sortie {proc.returncode})\n"
            f"{output}".rstrip())


# --------------------------------------------------------------------------
# Génération d'images (hors ligne) : Stable Diffusion si présent, sinon
# illustrations locales. Toujours honnête : jamais d'image "de remplissage".
# --------------------------------------------------------------------------

_VALID_OUT_EXT = {".png", ".jpg", ".jpeg", ".svg", ".pptx"}


def _resolve_sortie(sortie: str | None, folder: str, ext: str) -> str:
    d = os.path.join(config.WORKSPACE, folder)
    os.makedirs(d, exist_ok=True)
    if not sortie:
        return os.path.join(d, time.strftime("agent-%Y%m%d-%H%M%S") + ext)
    s = str(sortie).strip().strip('"')
    # Un dossier (ex: "images/") -> nom de fichier par défaut à l'intérieur.
    if s.endswith(("/", "\\")) or os.path.isdir(s):
        s = os.path.normpath(s)
        os.makedirs(s, exist_ok=True)
        return os.path.join(s, time.strftime("agent-%Y%m%d-%H%M%S") + ext)
    # Un nom de fichier sans extension -> on la rajoute.
    if os.path.splitext(s)[1].lower() not in _VALID_OUT_EXT:
        s += ext
    return os.path.abspath(s)


def _try_sd_api(prompt: str, sortie: str) -> str | None:
    """Moteur Stable Diffusion compatible API (ComfyUI / Automatic1111 / Forge)."""
    url = (config.SD_URL or "").strip()
    if not url:
        return None
    import requests as req
    endpoint = url.rstrip("/") + "/sdapi/v1/txt2img"
    try:
        if req.get(endpoint, timeout=4).status_code != 200:
            return None
    except Exception:
        return None
    try:
        r = req.post(endpoint, json={"prompt": prompt, "steps": 20,
                                     "width": 512, "height": 512, "seed": -1},
                     timeout=1200)
        images = r.json().get("images")
        if not images:
            return None
        import base64
        with open(sortie, "wb") as fh:
            fh.write(base64.b64decode(images[0]))
        return f"OK: image générée par Stable Diffusion dans {sortie}."
    except Exception:
        return None


def _try_sd_cpp(prompt: str, sortie: str) -> str | None:
    """Exécutable stable-diffusion.cpp + modèle .gguf/.safetensors configurés."""
    exe = (config.SDCPP or "").strip()
    model = (config.SD_MODEL or "").strip()
    if not exe or not model:
        return None
    if not (os.path.exists(exe) and os.path.exists(model)):
        return None
    try:
        proc = subprocess.run(
            [exe, "--model", model, "--prompt", prompt, "--output", sortie,
             "--height", "512", "--width", "512", "--steps", "20",
             "--cfg-scale", "5", "--seed", "-1"],
            capture_output=True, text=True, timeout=1200)
        if proc.returncode == 0 and os.path.exists(sortie):
            return f"OK: image générée par stable-diffusion.cpp dans {sortie}."
    except Exception:
        pass
    return None


def generer_image(prompt: str, sortie: str | None = None) -> str:
    """Génère une image hors ligne : SD local si présent, sinon illustration.

    Retourne TOUJOURS un message honnête : quel moteur a produit l'image, et
    le chemin du fichier créé (ou une erreur guidante si les données manquent).
    """
    prompt = (prompt or "").strip()
    if not prompt:
        return "ERREUR: description de l'image vide."
    out = _resolve_sortie(sortie, "images", ".png")
    os.makedirs(os.path.dirname(out), exist_ok=True)
    for engine in (_try_sd_api, _try_sd_cpp):
        msg = engine(prompt, out)
        if msg:
            return msg
    import dessin
    return dessin.generer(prompt, out)


# --------------------------------------------------------------------------
# Création de présentations PowerPoint (.pptx) pour enseignants
# (python-pptx, embarqué). Images réelles + liens vidéo cliquables.
# --------------------------------------------------------------------------

_PPTX_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff"}


def _add_pptx_bullets(tf, items: list) -> None:
    from pptx.util import Pt
    for idx, item in enumerate(items):
        item = str(item)
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        level = 1 if item.startswith("  ") else 0
        p.level = level
        p.text = item.strip().lstrip("-*•")
        p.font.size = Pt(22 if level == 0 else 18)


def creer_powerpoint(plan: str, sortie: str | None = None) -> str:
    """Crée un .pptx de cours à partir d'un plan JSON (slides, puces, notes,
    images existantes, liens vidéo). Ne crée jamais d'image de remplacement."""
    import json
    try:
        data = json.loads(plan)
    except Exception as err:
        return f"ERREUR: plan JSON invalide ({err}). Fournissez un objet JSON."
    if not isinstance(data, dict):
        return "ERREUR: le plan doit être un objet JSON {...}."
    slides = data.get("slides") or []
    if not isinstance(slides, list) or not slides:
        return "ERREUR: le plan doit contenir une liste non vide 'slides'."

    out = _resolve_sortie(sortie, "presentations", ".pptx")
    os.makedirs(os.path.dirname(out), exist_ok=True)

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        return ("ERREUR: python-pptx n'est pas installé. "
                "Installez-le avec: python -m pip install python-pptx.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    NAVY = RGBColor(15, 23, 42)
    BLUE = RGBColor(56, 189, 248)
    WHITE = RGBColor(255, 255, 255)
    TEXT_C = RGBColor(30, 41, 59)

    # --- Diapositive de titre -------------------------------------------------
    s = prs.slides.add_slide(blank)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                              prs.slide_width, Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    title = s.shapes.add_textbox(Inches(0.8), Inches(2.4),
                                 prs.slide_width - Inches(1.6), Inches(2))
    tf = title.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = str(data.get("titre") or "Présentation")
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    sub = (data.get("auteur") or "") + \
        ("   •   " + time.strftime("%d/%m/%Y") if data.get("auteur") else "")
    if sub:
        st = s.shapes.add_textbox(Inches(0.8), Inches(4.6),
                                  prs.slide_width - Inches(1.6), Inches(1))
        stf = st.text_frame
        stf.paragraphs[0].text = sub
        stf.paragraphs[0].font.size = Pt(20)
        stf.paragraphs[0].font.color.rgb = BLUE

    images_ok = 0
    videos_ok = 0
    images_skipped = []

    for num, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue
        s = prs.slides.add_slide(blank)
        # bandeau de titre
        band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  prs.slide_width, Inches(1.05))
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY
        band.line.fill.background()
        t = s.shapes.add_textbox(Inches(0.6), Inches(0.18),
                                 prs.slide_width - Inches(1.2), Inches(0.75))
        tf = t.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = str(slide.get("titre") or f"Diapositive {num}")
        tf.paragraphs[0].font.size = Pt(30)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE

        # texte + image à droite
        body = s.shapes.add_textbox(Inches(0.6), Inches(1.4),
                                    Inches(7.6), Inches(5.6))
        tf = body.text_frame
        tf.word_wrap = True
        _add_pptx_bullets(tf, slide.get("texte") or [])

        img_path = (slide.get("image") or "").strip()
        if img_path:
            real = _existing_path(img_path)
            ext = os.path.splitext(real)[1].lower()
            if os.path.isfile(real) and ext in _PPTX_IMAGE_EXTS:
                try:
                    s.shapes.add_picture(real, Inches(8.4), Inches(1.6),
                                         height=Inches(4.6))
                    images_ok += 1
                except Exception:
                    images_skipped.append(img_path)
            else:
                images_skipped.append(img_path)

        video = slide.get("video") or {}
        vurl = (video.get("url") or "").strip()
        if vurl:
            vtxt = (video.get("texte") or "Voir la vidéo").strip()
            box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.6), Inches(6.4), Inches(7.6), Inches(0.7))
            box.fill.solid()
            box.fill.fore_color.rgb = BLUE
            box.line.fill.background()
            vf = box.text_frame
            p = vf.paragraphs[0]
            run = p.add_run()
            run.text = "▶  " + vtxt
            run.font.size = Pt(18)
            run.font.color.rgb = NAVY
            run.hyperlink.address = vurl
            videos_ok += 1

        notes = (slide.get("notes") or "").strip()
        if notes:
            s.notes_slide.notes_text_frame.text = notes

    try:
        prs.save(out)
    except OSError as err:
        return f"ERREUR: impossible d'écrire le fichier: {err}"

    n = len(slides)
    msg = (f"OK: présentation créée dans {out} "
           f"({n} diapositive(s) de contenu + 1 titre, {images_ok} image(s) "
           f"intégrée(s), {videos_ok} lien(s) vidéo).")
    if images_skipped:
        msg += ("\nImage(s) NON intégrée(s) (fichier introuvable ou non image) : "
                + ", ".join(images_skipped))
    return msg


# --------------------------------------------------------------------------
# Lecture d'une IMAGE : modèle de vision Ollama si installé, sinon OCR
# (Tesseract), sinon un message qui explique comment activer la vision.
# --------------------------------------------------------------------------

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".bmp", ".gif", ".webp", ".tiff", ".jfif"}
# Ordre de préférence : les modèles récents lisent bien mieux le texte des
# captures d'écran que llava (qui hallucine). llava reste le dernier recours.
VISION_MODEL_HINTS = ("qwen2.5vl", "qwen2-vl", "llama3.2-vision", "gemma3",
                      "minicpm-v", "moondream", "llava")
_VISION_MAX_GB = 9.0  # au-delà, le modèle rame trop longtemps sur un CPU seul


def _available_vision_model() -> str | None:
    """Trouve un modèle de vision déjà téléchargé dans Ollama (ou None).

    Un modèle forcé par la variable d'environnement AGENT_VISION_MODEL est
    prioritaire (s'il est bien installé) ; sinon on cherche par indice de nom
    ou par la capacité "vision" annoncée par Ollama, en préférant un modèle
    assez léger pour tourner sur CPU. On ne renvoie JAMAIS un modèle sans
    capacité vision (l'envoi d'images échouerait).
    """
    import requests as req
    try:
        data = req.get(f"{config.OLLAMA_URL}/api/tags", timeout=5).json()
    except Exception:
        return None
    installed = []
    for m in data.get("models", []):
        caps = m.get("capabilities") or []
        name = m.get("name", "")
        installed.append({
            "name": name, "size": m.get("size") or 0,
            "vision": "vision" in caps or _hint_for(name) is not None,
        })
    forced = (config.VISION_MODEL or "").strip()
    if forced:
        for mod in installed:
            if mod["name"] == forced or mod["name"].startswith(forced + ":"):
                return mod["name"]
    fallback = None
    for hint in VISION_MODEL_HINTS:
        match = None
        for mod in installed:
            if mod["name"].startswith(hint):
                if match is None or mod["size"] < match[1]:
                    match = (mod["name"], mod["size"])
        if match:
            if match[1] <= _VISION_MAX_GB * 1024 ** 3:
                return match[0]
            fallback = fallback or match[0]   # trop lourd, gardé en secours
    # Dernier recours : un modèle qui annonce explicitement la capacité vision,
    # même s'il n'est pas dans nos indices de noms (nouveau modèle, renommé...).
    light = None
    for mod in installed:
        if not mod["vision"] or mod["name"].startswith(tuple(VISION_MODEL_HINTS)):
            continue
        if light is None or mod["size"] < light[1]:
            light = (mod["name"], mod["size"])
    if light and light[1] <= _VISION_MAX_GB * 1024 ** 3:
        return light[0]
    return fallback or None


def _hint_for(name: str) -> str | None:
    """Indice de modèle vision dont le nom commence par ``name``, sinon None."""
    for hint in VISION_MODEL_HINTS:
        if name.startswith(hint):
            return hint
    return None


def _image_tiles(path: str, max_w: int = 1024, max_h: int = 768) -> list:
    """Découpe une grande image en tuiles lisibles pour le modèle de vision.

    Les petits modèles de vision (llava...) réduisent l'image à ~336 px :
    sur une capture d'écran large, le texte devient illisible et le modèle
    INVENTE du contenu. Découper en tuiles (avec chevauchement) garde un
    texte réellement lisible.
    """
    from PIL import Image
    im = Image.open(path)
    if im.mode != "RGB":
        im = im.convert("RGB")
    w, h = im.size
    if max(w, h) < 320:                     # toute petite image : on agrandit
        scale = 480 / max(w, h)
        im = im.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
        w, h = im.size
    if w <= max_w and h <= max_h:
        return [im]
    cols = max(1, math.ceil(w / max_w))
    rows = max(1, math.ceil(h / max_h))
    tiles = []
    for r in range(rows):
        for c in range(cols):
            x0 = max(0, int(w * c / cols) - 16)
            x1 = min(w, int(w * (c + 1) / cols) + 16)
            y0 = max(0, int(h * r / rows) - 16)
            y1 = min(h, int(h * (r + 1) / rows) + 16)
            tiles.append(im.crop((x0, y0, x1, y1)))
    return tiles


_VISION_PROMPT = (
    "Tu es un lecteur de captures d'écran précis et honnête.\n"
    "1. Décris en UNE phrase ce que montre l'image.\n"
    "2. Transcris mot pour mot TOUT le texte visible (questions, réponses, "
    "boutons, titres), en gardant la structure.\n"
    "RÈGLE ABSOLUE : n'invente JAMAIS de phrase, numéro, question ou réponse. "
    "Si un texte n'est pas parfaitement lisible, écris [texte illisible] à sa "
    "place. Si l'image ne contient aucun texte, dis simplement : Aucun texte visible."
)


def _vision_describe(path: str, model: str) -> str:
    """Envoie l'image (éventuellement découpée en tuiles) au modèle de vision."""
    import base64
    import io
    import requests as req
    tiles = _image_tiles(path)
    parts: list[str] = []
    errors: list[str] = []
    for i, tile in enumerate(tiles, start=1):
        buf = io.BytesIO()
        tile.save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode("ascii")
        content = _VISION_PROMPT
        if len(tiles) > 1:
            content += f"\n(partie {i}/{len(tiles)} de l'image)"
        payload = {
            "model": model,
            "messages": [{"role": "user", "content": content,
                          "images": [b64]}],
            "stream": False,
            "options": {"num_ctx": 2048, "temperature": 0.1},
        }
        try:
            resp = req.post(f"{config.OLLAMA_URL}/api/chat",
                            json=payload, timeout=180)
        except Exception as err:
            errors.append(f"'{model}' n'a pas répondu: {err}")
            continue
        if resp.status_code != 200:
            errors.append(f"'{model}' (HTTP {resp.status_code}): {resp.text[:300]}")
            continue
        resp.encoding = "utf-8"
        text = (resp.json().get("message", {}).get("content") or "").strip()
        if not text:
            errors.append(f"'{model}' n'a rien retourné.")
            continue
        if len(tiles) > 1:
            parts.append(f"[Partie {i}/{len(tiles)} de l'image]\n{text}")
        else:
            parts.append(text)
    if errors and not parts:
        return errors[0]
    body = "\n\n".join(parts)
    suffix = (f"\n\n(Attention : {'; '.join(errors)})" if errors else "")
    return (f"Description de l'image {os.path.basename(path)} (modèle {model}) :\n"
            f"{body}{suffix}")


def _ocr_text(path: str) -> str:
    """Extrait le texte d'une image via Tesseract, si installé (optionnel)."""
    try:
        from PIL import Image
        import pytesseract
    except ImportError:
        return ""
    try:
        text = pytesseract.image_to_string(Image.open(path), lang="fra+eng")
    except Exception:
        return ""
    text = (text or "").strip()
    if not text:
        return ""
    return (f"Texte extrait de l'image {os.path.basename(path)} par OCR "
            f"(Tesseract) :\n{text[:DOC_TEXT_LIMIT]}")


# OCR NATIF WINDOWS (Windows.Media.Ocr) via Windows PowerShell 5.1 : exact,
# gratuit, 100 % hors ligne, présent sur tout Windows 10/11. Indispensable
# pour LIRE le texte des captures d'écran : un modèle de vision local comme
# llava INVENTE du contenu quand le texte devient trop petit à ~336 px.
_POWERSHELL_OCR_SCRIPT = r"""
[Console]::OutputEncoding = [System.Text.Encoding]::UTF8
Add-Type -AssemblyName System.Runtime.WindowsRuntime
$null = [Windows.Media.Ocr.OcrEngine,Windows.Foundation,ContentType=WindowsRuntime]
$null = [Windows.Graphics.Imaging.BitmapDecoder,Windows.Foundation,ContentType=WindowsRuntime]
$null = [Windows.Storage.StorageFile,Windows.Foundation,ContentType=WindowsRuntime]
function Await($WinRtTask, $ResultType) {
    $asTask = ([System.WindowsRuntimeSystemExtensions].GetMethods() |
        Where-Object { $_.Name -eq 'AsTask' -and $_.GetParameters().Count -eq 1 -and
                       $_.GetParameters()[0].ParameterType.Name -eq 'IAsyncOperation`1' })[0]
    $netTask = $asTask.MakeGenericMethod($ResultType).Invoke($null, @($WinRtTask))
    $netTask.Wait(-1) | Out-Null
    $netTask.Result
}
$path = $args[0]
try {
    $file = Await ([Windows.Storage.StorageFile]::GetFileFromPathAsync($path)) ([Windows.Storage.StorageFile])
    $stream = Await ($file.OpenAsync([Windows.Storage.FileAccessMode]::Read)) ([Windows.Storage.Streams.IRandomAccessStream])
    $decoder = Await ([Windows.Graphics.Imaging.BitmapDecoder]::CreateAsync($stream)) ([Windows.Graphics.Imaging.BitmapDecoder])
    $bitmap = Await ($decoder.GetSoftwareBitmapAsync()) ([Windows.Graphics.Imaging.SoftwareBitmap])
    $engine = [Windows.Media.Ocr.OcrEngine]::TryCreateFromUserProfileLanguages()
    if (-not $engine) { Write-Output '__NO_OCR_ENGINE__'; exit }
    $result = Await ($engine.RecognizeAsync($bitmap)) ([Windows.Media.Ocr.OcrResult])
    foreach ($line in $result.Lines) { Write-Output $line.Text }
} catch {
    Write-Output '__OCR_ERROR__'
}
"""


def _windows_ocr(path: str) -> str:
    """OCR natif Windows : texte EXACT d'une capture d'écran.

    Renvoie '' si indisponible (pas de Windows PowerShell 5.1, pas de module
    de langue OCR, ou échec). Le texte renvoyé fait foi : il n'est pas
    "interprété" par un modèle.
    """
    import tempfile
    ps = shutil.which("powershell")      # Windows PowerShell 5.1 (pas pwsh)
    if not ps:
        return ""
    path = os.path.abspath(path)         # WinRT exige un chemin ABSOLU
    fd, script = tempfile.mkstemp(suffix=".ps1")
    try:
        with os.fdopen(fd, "w", encoding="utf-8") as fh:
            fh.write(_POWERSHELL_OCR_SCRIPT)
        try:
            r = subprocess.run(
                [ps, "-NoProfile", "-NonInteractive", "-ExecutionPolicy",
                 "Bypass", "-File", script, path],
                capture_output=True, timeout=120)
        except Exception:
            return ""
        if r.returncode != 0:
            return ""
        out = r.stdout.decode("utf-8", errors="replace")
    finally:
        try:
            os.remove(script)
        except Exception:
            pass
    lines = [ln.strip() for ln in out.splitlines()
             if ln.strip() and ln.strip() not in ("__NO_OCR_ENGINE__",
                                                  "__OCR_ERROR__")]
    if not lines:
        return ""
    return (f"Texte extrait de l'image {os.path.basename(path)} "
            f"par OCR Windows (exact) :\n" + "\n".join(lines[:80]))


def read_image(path: str) -> str:
    """Décrit une image (contenu + texte visible). 100 % hors ligne :

      1. OCR Windows natif (exact) puis Tesseract -> le TEXTE fait foi ;
      2. modèle de vision local (llava...) -> description de ce que montre
         l'image (photos, schémas) ;
      3. sinon, un message qui explique comment activer la vision.
    """
    path = _existing_path(path)
    if not os.path.isfile(path):
        return f"ERREUR: image introuvable: {path}"
    if os.path.splitext(path)[1].lower() not in IMAGE_EXTS:
        return (f"ERREUR: '{path}' n'est pas une image reconnue "
                f"({', '.join(sorted(IMAGE_EXTS))}).")
    size = os.path.getsize(path)
    if size > 10 * 1024 * 1024:
        return "ERREUR: image trop volumineuse (> 10 Mo)."

    parts: list[str] = []
    # 1) Texte EXACT (OCR Windows natif, sinon Tesseract). C'est lui qui fait
    #    foi pour lire une question, un exercice, une capture d'écran.
    ocr = _windows_ocr(path) or _ocr_text(path)
    if ocr:
        body = ocr.split("\n", 1)[1] if "\n" in ocr else ocr
        if len(body) >= 40:
            # Une capture d'écran avec du texte : l'OCR suffit, exact et rapide.
            # Inutile (et risqué) de laisser llava réécrire le texte de travers.
            return ocr
        parts.append(ocr)
    # 2) Description du modèle de vision (approximative, complémentaire) :
    #    utile pour une photo ou un schéma, où l'OCR ne trouve rien.
    vision = _available_vision_model()
    if vision:
        parts.append(_vision_describe(path, vision))
    if parts:
        return "\n\n".join(parts)
    return (f"Le fichier '{path}' est une image ({size:,} o). Impossible de la "
            f"DÉCRIRE hors ligne : aucun modèle de vision ni OCR installé.\n"
            f"Pour activer la vision, installez l'un des deux :\n"
            f"  1. ollama pull llava:7b   (modèle de vision local)\n"
            f"  2. Tesseract OCR + pytesseract + Pillow (texte des images)\n"
            f"En attendant, je peux résumer le fichier par son nom ou lire "
            f"tout document avec read_document.")


# --------------------------------------------------------------------------
# Calcul symbolique VÉRIFIÉ (SymPy).
#
# Le modèle NE DOIT JAMAIS écrire le code SymPy à la main (il fait des fautes
# de transcription : `sympy.ln(x)+1` au lieu de `sympy.ln(x+1)`). Il passe
# l'expression EN TEXTE BRUT, copiée telle quelle, et SymPy se charge de
# l'interpréter, de calculer ET de vérifier par dérivation.
# --------------------------------------------------------------------------

# Noms de fonctions connus : on n'y insère pas de multiplication implicite.
_FUNC_NAMES = {
    "ln", "log", "exp", "sin", "cos", "tan", "asin", "acos", "atan",
    "sinh", "cosh", "tanh", "sqrt", "abs", "erf", "sec", "csc", "cot",
}


def _norm_op(operation: str) -> str:
    """Normalise le nom de l'opération demandée ('intégrale', 'dérivée'...)."""
    op = (operation or "").lower()
    for a, b in (("é", "e"), ("è", "e"), ("ê", "e"), ("à", "a"),
                 ("ô", "o"), ("î", "i"), ("û", "u"), ("ç", "c")):
        op = op.replace(a, b)
    if not op:
        return "integrale"
    if op.startswith("integr") or op.startswith("primitive"):
        return "integrale"
    if op.startswith("deriv"):
        return "derivee"
    if op.startswith("equat") or op.startswith("resoud") or op.startswith("resol"):
        return "equation"
    if op.startswith("simpl"):
        return "simplifier"
    if op.startswith("factor"):
        return "factoriser"
    if op.startswith("limit"):
        return "limite"
    return "integrale"


def _prepare_expr(s: str) -> str:
    """Petits accommodements avant l'interprétation SymPy : x^2 -> x**2,
    multiplication implicite (2x -> 2*x, x(x+1) -> x*(x+1))."""
    s = s.replace("^", "**").replace("√", "sqrt").strip()
    # ')' suivi d'une lettre/chiffre/'(' -> multiplication.
    s = re.sub(r"\)\s*(?=[A-Za-z0-9(])", ")*", s)
    # chiffre suivi d'une lettre ou d'une parenthèse -> multiplication.
    s = re.sub(r"(?<=\d)\s*(?=[A-Za-z(])", "*", s)
    # variable simple (x, y, t, z) directement suivie de '(' -> multiplication,
    # sauf si elle fait partie d'un nom de fonction (sin(, ln(, exp(...).
    s = re.sub(r"(^|[^\w])([xytz])\s*\(", r"\1\2*(", s)
    return s


def _by_parts_steps(f, sym, sympy):
    """Dérivation par parties VÉRIFIÉE pour f = ln(a*x + b).

    Chaque étape est recalculée par SymPy : la dérivation est donc exacte.
    Renvoie '' si la fonction n'est pas de la forme ln(linéaire).
    """
    if not isinstance(f, sympy.log):
        return ""
    g = f.args[0]
    if sympy.degree(g, sym) != 1:
        return ""
    p = sympy.Poly(g, sym)
    a, b = p.nth(1), p.nth(0)
    du = sympy.simplify(a / g)
    v_int = sympy.simplify(sympy.integrate(sym * a / g, sym))
    final = sympy.simplify(sym * f - v_int)
    g_s = str(g).replace("log(", "ln(")
    du_s = str(du).replace("log(", "ln(")
    v_s = str(v_int).replace("log(", "ln(")
    f_s = str(final).replace("log(", "ln(")
    return (f"\n\nMÉTHODE PAR PARTIES (chaque étape vérifiée par SymPy) :\n"
            f"u = ln({g_s}),  dv = d{sym}\n"
            f"du = {du_s} d{sym},  v = {sym}\n"
            f"∫ u·dv = u·v − ∫ v·du = {sym}·ln({g_s}) − ∫ {sym}·{a}/({g_s}) d{sym}\n"
            f"∫ {sym}·{a}/({g_s}) d{sym} = {v_s}\n"
            f"Résultat : {sym}·ln({g_s}) − ({v_s}) = {f_s} + C")


def _strip_integral(s: str) -> tuple[str, str | None]:
    """Retire le symbole '∫' éventuel et la différentielle finale.

    Les étudiants écrivent '∫ln(x+1)dx' ou '∫ ln(x+1) dx' : ces notations
    doivent devenir l'expression propre 'ln(x+1)'. Retourne (expression
    propre, variable inférée depuis la différentielle ou None).
    """
    s = s.replace("∫", "").replace("∬", "").strip()
    m = re.search(r"d\s*\(\s*([A-Za-z])\s*\)\s*$", s)   # d(x), d x entre parenthèses
    if m:
        return s[:m.start()].strip(), m.group(1)
    m = re.search(r"\bd\s*([A-Za-z])\s*$", s)           # dx, d x, dx final
    if m:
        return s[:m.start()].strip(), m.group(1)
    return s, None


def calcul_symbolique(expression: str, operation: str = "integrale",
                      variable: str = "x", point: str = "oo") -> str:
    """Calcule et VÉRIFIE une expression mathématique avec SymPy.

    `expression` est fournie en texte brut (ex: 'ln(x+1)'), copiée telle
    quelle de la demande de l'étudiant. SymPy l'interprète, calcule, puis
    vérifie le résultat (la dérivée du résultat doit redonner l'expression).
    """
    try:
        import sympy
    except ImportError:
        return ("ERREUR: SymPy n'est pas installé. "
                "Installez-le avec: python -m pip install sympy.")
    op = _norm_op(operation)
    expr = (expression or "").strip()
    if not expr:
        return ("ERREUR: expression vide. Passez SEULEMENT la fonction, "
                "ex: expression='ln(x+1)'.")
    # Nettoie la notation '∫ln(x+1)dx' -> 'ln(x+1)' (et infère la variable).
    expr, var_inf = _strip_integral(expr)
    if var_inf:
        variable = var_inf
    if not expr:
        return "ERREUR: expression vide après nettoyage. Passez la fonction, ex: 'ln(x+1)'."
    sym = sympy.symbols(variable)
    locals_map = {variable: sym, "ln": sympy.log, "log": sympy.log,
                  "e": sympy.E, "pi": sympy.pi, "oo": sympy.oo}
    try:
        if "=" in expr:
            lhs_s, _, rhs_s = expr.partition("=")
            f = sympy.sympify(_prepare_expr(lhs_s), locals=locals_map) - \
                sympy.sympify(_prepare_expr(rhs_s), locals=locals_map)
        else:
            f = sympy.sympify(_prepare_expr(expr), locals=locals_map)
    except Exception as err:
        return (f"ERREUR: expression '{expression}' non comprise ({err}). "
                f"Passez SEULEMENT la fonction en texte brut, ex: 'ln(x+1)', "
                f"'x**2*exp(x)' ou l'équation 'x^2 - 5*x + 6 = 0'.")
    f = sympy.simplify(f)
    # Garde anti-interprétation sauvage : si '∫', 'dx' ou une intégrale non
    # résolue traînent encore, l'expression n'a pas été comprise correctement.
    if ("∫" in str(f) or re.search(r"\bdx\b", str(f))
            or "Integral(" in str(f)):
        return (f"ERREUR: expression '{expression}' mal interprétée (SymPy a "
                f"lu '{f}'). Passez la fonction sans le symbole ∫ ni le dx, "
                f"ex: expression='ln(x+1)'.")

    if op == "derivee":
        r = sympy.simplify(sympy.diff(f, sym))
        check = sympy.simplify(sympy.integrate(r, sym) - f)
        verif = "CORRECT" if check == 0 else f"à vérifier: reste {check}"
        return (f"Expression interprétée : {f}\n"
                f"f'({variable}) = {r}\n"
                f"VÉRIFICATION : {verif}")

    if op == "equation":
        try:
            sols = sympy.solve(f, sym)
        except Exception as err:
            return f"ERREUR: résolution impossible: {err}"
        # Vérifie chaque solution en la substituant dans l'équation.
        try:
            checks = [sympy.simplify(f.subs(sym, s)) for s in sols]
            verif = "CORRECT" if all(c == 0 for c in checks) else \
                f"DOUTEUX: restes {[c for c in checks if c != 0]}"
        except Exception:
            verif = "calculée (vérification par substitution impossible)"
        return (f"Équation : {expr}\nSolution(s) {variable} : {sols}\n"
                f"VÉRIFICATION : {verif}")

    if op == "limite":
        try:
            p = sympy.sympify(_prepare_expr(point), locals=locals_map)
            r = sympy.limit(f, sym, p)
        except Exception as err:
            return f"ERREUR: limite impossible: {err}"
        return f"Expression interprétée : {f}\nlim {variable}→{point} : {r}"

    if op == "simplifier":
        return f"Expression interprétée : {expr}\nSimplifiée : {sympy.simplify(f)}"

    if op == "factoriser":
        try:
            return f"Expression interprétée : {expr}\nFactorisée : {sympy.factor(f)}"
        except Exception as err:
            return f"ERREUR: factorisation impossible: {err}"

    # Intégrale (cas par défaut aussi).
    r = sympy.integrate(f, sym)
    if "Integral(" in str(r):
        return (f"ERREUR: SymPy n'a pas su intégrer '{expr}'. "
                f"Vérifiez l'expression ou simplifiez-la.")
    check = sympy.simplify(sympy.diff(r, sym) - f)
    verif = "CORRECT (la dérivée du résultat redonne exactement f(x))" \
        if check == 0 else f"DOUTEUX: il reste {check}"
    form = ""
    try:
        fr = sympy.collect(r, sympy.log(sym + 1))
        if fr != r:
            form = f"\nForme regroupée : {fr} + C"
    except Exception:
        pass
    try:
        by_parts = _by_parts_steps(f, sym, sympy)
    except Exception:
        by_parts = ""
    return (f"Expression interprétée : {f}\n"
            f"∫ f({variable}) d{variable} = {r} + C{form}"
            f"{by_parts}\n"
            f"VÉRIFICATION : {verif}")


# --------------------------------------------------------------------------
# Recherche d'un terme dans un dossier (type "grep").
# --------------------------------------------------------------------------

# Dossiers à ignorer lors d'une recherche récursive.
IGNORED_DIRS = {"node_modules", ".git", "dist", ".next", "coverage", "__pycache__"}


def count_occurrences(path: str, term: str) -> str:
    """Compte EXACTEMENT les occurrences d'un mot (mots entiers, insensible
    à la casse). Décompte par le code : résultat fiable à 100 %."""
    if not os.path.isfile(path):
        return f"ERREUR: fichier introuvable: {path}"
    try:
        with open(path, encoding="utf-8", errors="replace") as fh:
            text = fh.read()
    except OSError as err:
        return f"ERREUR: {err}"
    if len(text) > 5 * 1024 * 1024:
        return f"ERREUR: fichier trop volumineux ({len(text)} octets)."
    if " " in term.strip():
        count = text.lower().count(term.lower())
    else:
        count = len(re.findall(rf"\b{re.escape(term)}\b", text,
                               flags=re.IGNORECASE))
    return (f"Le terme '{term}' apparaît {count} fois dans {path} "
            f"({len(text.split())} mots dans le fichier).")


def search_in_files(term: str, folder: str = ".") -> str:
    """Recherche un mot/expression (insensible à la casse) dans un dossier.

    Retourne `fichier:n°ligne: contenu`, comme le font les bons outils.
    Limite à 100 résultats pour ne pas saturer le contexte du modèle.
    """
    needle = term.lower()
    hits: list[str] = []
    for root, dirs, files in os.walk(folder):
        dirs[:] = [d for d in dirs if d not in IGNORED_DIRS]   # on évite les dossiers lourds
        for name in files:
            full = os.path.join(root, name)
            try:
                with open(full, "rb") as fh:
                    data = fh.read()
            except OSError:
                continue                                      # fichier illisible : on ignore
            if b"\0" in data[:8192]:                          # binaire : on saute
                continue
            for i, line in enumerate(data.decode("utf-8", errors="replace").split("\n"),
                                     start=1):
                if needle in line.lower():
                    hits.append(f"{full}:{i}: {line.strip()[:200]}")
                    if len(hits) >= 100:
                        break
            if len(hits) >= 100:
                break
        if len(hits) >= 100:
            break
    if not hits:
        return f"Aucune occurrence de '{term}' dans {folder}."
    return f"{len(hits)} occurrence(s) de '{term}' :\n" + "\n".join(hits)


# --------------------------------------------------------------------------
# Lecture de documents : un .docx/.xlsx/.pptx est un ZIP (zipfile, natif),
# un .pdf contient des flux texte compressés (zlib, natif). Zéro dépendance.
# Limite connue : un PDF numérisé (image) n'a pas de texte à extraire.
# --------------------------------------------------------------------------

# Taille max du texte renvoyé au modèle (on protège son contexte).
DOC_TEXT_LIMIT = 8000

# Ponctuation française des PDF encodés en CP1252 (les apostrophes, tirets…).
_CP1252 = {
    0x80: "€", 0x82: "‚", 0x83: "ƒ", 0x84: "„", 0x85: "…", 0x86: "†", 0x87: "‡",
    0x88: "ˆ", 0x89: "‰", 0x8A: "Š", 0x8B: "‹", 0x8C: "Œ", 0x8E: "Ž",
    0x91: "‘", 0x92: "’", 0x93: "“", 0x94: "”", 0x95: "•", 0x96: "–", 0x97: "—",
    0x98: "˜", 0x99: "™", 0x9A: "š", 0x9B: "›", 0x9C: "œ", 0x9E: "ž", 0x9F: "Ÿ",
}


def _decode_xml(s: str) -> str:
    """Décode les entités XML (&amp;, &lt;…) en caractères réels."""
    return (s.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
            .replace("&quot;", '"').replace("&apos;", "'"))


def _xml_to_text(xml: str) -> str:
    """Transforme du XML Word/PowerPoint en texte lisible (balises retirées)."""
    text = re.sub(r"</w:p>|</a:p>", "\n", xml)      # fin de paragraphe -> ligne
    text = re.sub(r"<w:tab[^>]*/>|<a:br[^>]*/>", "\t", text)
    text = re.sub(r"<[^>]+>", "", text)             # toutes les autres balises
    text = _decode_xml(text)
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def _docx_to_text(path: str) -> str:
    """Word : le texte vit dans word/document.xml (un XML aux balises <w:…>)."""
    with zipfile.ZipFile(path) as zf:
        xml = zf.read("word/document.xml").decode("utf-8")
    return _xml_to_text(xml)


def _pptx_to_text(path: str) -> str:
    """PowerPoint : chaque diapositive est un XML dans ppt/slides/slideN.xml."""
    with zipfile.ZipFile(path) as zf:
        slides = sorted(
            (n for n in zf.namelist() if re.match(r"^ppt/slides/slide\d+\.xml$", n)),
            key=lambda n: int(re.search(r"\d+", n).group()),
        )
        if not slides:
            return "(aucune diapositive trouvée)"
        parts = []
        for name in slides:
            num = int(re.search(r"\d+", name).group())
            xml = zf.read(name).decode("utf-8")
            parts.append(f"Diapositive {num}:\n{_xml_to_text(xml)}")
    return "\n\n".join(parts)


def _xlsx_shared_strings(xml: str) -> list[str]:
    """Excel : les textes partagés sont dans xl/sharedStrings.xml."""
    strings: list[str] = []
    for si in re.findall(r"<si\b[\s\S]*?</si>", xml):
        text = "".join(re.findall(r"<t[^>]*>([\s\S]*?)</t>", si))
        strings.append(_decode_xml(text))
    return strings


def _xlsx_sheet(xml: str, shared: list[str]) -> str:
    """Une feuille : chaque cellule <c> devient `référence: valeur`."""
    lines: list[str] = []
    for cell in re.findall(r"<c\b[^>]*>[\s\S]*?</c>", xml):
        ref = re.search(r'r="([A-Z]+\d+)"', cell)
        if not ref:
            continue
        ref = ref.group(1)
        typ = re.search(r't="([^"]*)"', cell)
        typ = typ.group(1) if typ else ""
        v = re.search(r"<v>([\s\S]*?)</v>", cell)
        v = v.group(1) if v else None
        inline = re.search(r"<t[^>]*>([\s\S]*?)</t>", cell)
        inline = inline.group(1) if inline else None
        if typ == "s" and v is not None:
            value = shared[int(v)] if int(v) < len(shared) else ""
        elif inline is not None:
            value = inline
        elif v is not None:
            value = v
        else:
            value = ""
        if value.strip():
            lines.append(f"{ref}: {_decode_xml(value)}")
    return "\n".join(lines)


def _xlsx_to_text(path: str) -> str:
    """Excel : feuilles dans l'ordre, cellules avec leurs références (A1, B2…)."""
    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()
        if "xl/sharedStrings.xml" in names:
            shared = _xlsx_shared_strings(zf.read("xl/sharedStrings.xml").decode("utf-8"))
        else:
            shared = []
        sheets = sorted(
            (n for n in names if re.match(r"^xl/worksheets/sheet\d+\.xml$", n)),
            key=lambda n: int(re.search(r"\d+", n).group()),
        )
        if not sheets:
            return "(aucune feuille trouvée)"
        parts = []
        for i, name in enumerate(sheets, start=1):
            xml = zf.read(name).decode("utf-8")
            parts.append(f"Feuille {i}:\n{_xlsx_sheet(xml, shared)}")
    return "\n\n".join(parts)


def _pdf_streams(data: bytes) -> list[bytes]:
    """Décompresse les flux texte (FlateDecode) d'un PDF, avec ou sans en-tête zlib."""
    streams: list[bytes] = []
    for m in re.finditer(rb"stream\r?\n(.*?)\r?\nendstream", data, re.DOTALL):
        raw = m.group(1)
        for wbits in (15, -15):                     # -15 = deflate brut, 15 = zlib complet
            try:
                streams.append(zlib.decompress(raw, wbits))
                break
            except zlib.error:
                continue
    return streams


def _decode_pdf_string(s: str) -> str:
    """Décode une chaîne PDF : échappements + encodage CP1252."""
    s = s.replace(r"\(", "(").replace(r"\)", ")").replace(r"\\", "\\")
    s = s.replace(r"\n", "\n").replace(r"\r", "")
    s = re.sub(r"\\([0-7]{1,3})", lambda m: chr(int(m.group(1), 8)), s)
    return "".join(_CP1252.get(ord(ch), ch) if "\x80" <= ch <= "\x9f" else ch
                   for ch in s)


def _pdf_op_text(op: str) -> str:
    """Texte d'un opérateur Tj ou TJ (les parenthèses délimitantes sont retirées)."""
    if op.startswith("("):
        return _decode_pdf_string(op[1:op.rfind(")")])
    arr = op[1:op.rfind("]")]                       # le tableau de chaînes de TJ
    return "".join(_decode_pdf_string(p)
                   for p in re.findall(r"\(((?:[^()\\]|\\.)*)\)", arr))


def _extract_pdf_text(data: bytes) -> str:
    """Recolle les morceaux de texte d'un PDF (les mots sont souvent découpés)."""
    lines: list[str] = []
    for stream in _pdf_streams(data):
        chunk = stream.decode("latin1")
        for segment in re.split(r"T\*|Td|TD", chunk):
            ops = re.findall(
                r"\((?:[^()\\]|\\.)*\)\s*Tj|\[(?:[^\[\]\\]|\\.)*\]\s*TJ",
                segment)
            if not ops:
                continue
            line = "".join(_pdf_op_text(op) for op in ops)
            line = re.sub(r"\s{2,}", " ", line).strip()
            if line:
                lines.append(line)
    return "\n".join(lines).strip()


def read_document(path: str) -> str:
    """Extrait le texte lisible d'un document : Word, PowerPoint, Excel, PDF, texte."""
    path = _existing_path(path)
    if not os.path.isfile(path):
        parent = os.path.dirname(os.path.abspath(path))
        try:
            listing = "\n".join(
                f"  - {name}" for name in sorted(os.listdir(parent))[:40])
        except OSError:
            listing = ""
        hint = (f"\nListe du dossier {parent}:\n{listing}" if listing else "")
        return (f"ERREUR: fichier introuvable: {path}"
                f"\nVérifiez le chemin ou cherchez le fichier.{hint}")
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    try:
        if ext == "docx":
            text = _docx_to_text(path)
        elif ext == "pptx":
            text = _pptx_to_text(path)
        elif ext == "xlsx":
            text = _xlsx_to_text(path)
        elif ext == "pdf":
            with open(path, "rb") as fh:
                text = _extract_pdf_text(fh.read())
        else:                                       # fichier texte ordinaire
            with open(path, "rb") as fh:
                data = fh.read()
            if b"\0" in data[:8192]:
                return (f"ERREUR: type de fichier '{ext}' non supporté par "
                        f"read_document, et le fichier semble binaire.")
            text = data.decode("utf-8", errors="replace")
    except (OSError, zipfile.BadZipFile) as err:
        return f"ERREUR: impossible d'extraire le texte: {err}"
    text = (text or "").strip()
    if not text:
        return "(document vide ou sans texte lisible)"
    if len(text) > DOC_TEXT_LIMIT:
        return text[:DOC_TEXT_LIMIT] + "\n...[tronqué]"
    return text


# Correspondance nom d'outil -> fonction Python.
EXECUTORS = {
    "list_dir": list_dir,
    "read_file": read_file,
    "read_document": read_document,
    "search_in_files": search_in_files,
    "count_occurrences": count_occurrences,
    "write_file": write_file,
    "edit_file": edit_file,
    "glob": glob,
    "bash": bash,
    "calcul_symbolique": calcul_symbolique,
    "read_image": read_image,
    "generer_image": generer_image,
    "creer_powerpoint": creer_powerpoint,
}


def execute_tool(name: str, arguments: dict) -> str:
    """Point d'entrée unique : exécute un outil par son nom, en sécurité."""
    if name not in EXECUTORS:
        return f"ERREUR: outil inconnu '{name}'."
    try:
        return EXECUTORS[name](**arguments)
    except TypeError as err:
        return f"ERREUR: mauvais arguments pour {name}: {err}"
    except Exception as err:                     # pragma: no cover - filet de sécurité
        return f"ERREUR inattendue dans {name}: {err}"
